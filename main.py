import streamlit as st
import base64
import groq
import pptx
from pptx.util import Inches, Pt
import os
from dotenv import load_dotenv

load_dotenv()

GROQ_API_KEY = os.getenv("GROQ_API_KEY")

client = groq.Client(api_key=GROQ_API_KEY)

# Custom formatting options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

def generate_slide_titles(topic):
    prompt = f"Generate 5 Slide titles for given topic '{topic}'."
    response = client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=200
    )
    
    # Access content from the response
    message_content = response.choices[0].message.content
    print(message_content)
    return message_content.split("\n")


def generate_slide_content(slide_title, topic):
    prompt = f"Generate Content for the Slide Title {slide_title} for given topic {topic}."
    response = client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=1000
    )
    
    # Access content from the response
    message_content = response.choices[0].message.content
    print(message_content)
    return message_content

def create_presentation(topic, slide_titles, slide_contents):
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[1]
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide.shapes.placeholders[1].text = slide_content

        # Customized font sizes
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE
        
    prs.save(f"generated_ppt/{topic}_presentation.pptx")
