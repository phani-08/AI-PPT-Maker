import streamlit as st
import base64
import groq
import pptx
from pptx.util import Inches, Pt
import os
from dotenv import load_dotenv

load_dotenv()

GROQ_API_KEY = os.getenv("GROQ_API_KEY")

client = groq.Client(api_key=GROQ_API_KEY)

# Custom formatting options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

def generate_slide_titles(topic):
    prompt = f"Generate 5 Slide titles for given topic '{topic}'."
    response = client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=200
    )
    
    # Access content from the response
    message_content = response.choices[0].message.content
    print(message_content)
    return message_content.split("\n")


def generate_slide_content(slide_title, topic):
    prompt = f"Generate Content for the Slide Title {slide_title} for given topic {topic}. Remember it should not be so long exceeding slide better concise info which fits to a slide"
    response = client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=300
    )
    
    # Access content from the response
    message_content = response.choices[0].message.content
    print(message_content)
    return message_content

def create_presentation(topic, slide_titles, slide_contents):
    prs = pptx.Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Set the slide title
        title_shape = slide.shapes.title
        title_shape.text = slide_title
        title_shape.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE

        # Use the body placeholder for content
        body_shape = slide.shapes.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()  # Remove any existing paragraphs

        # Process markdown content: split content into lines
        for line in slide_content.split("\n"):
            # Create a new paragraph for each line
            p = text_frame.add_paragraph()
            # Check for markdown bullet points
            if line.strip().startswith("- "):
                p.text = line.strip()[2:]  # Remove the bullet marker
                p.level = 1                # Set the bullet indentation level
                p.font.size = SLIDE_FONT_SIZE
                p.font.name = "Arial"
                p.bullet = True
            # Check for markdown headings (e.g., starting with '#' for a heading)
            elif line.strip().startswith("#"):
                p.text = line.strip().lstrip("#").strip()
                p.font.size = Pt(24)       # Apply a larger font size for headings
                p.font.bold = True
            else:
                p.text = line.strip()
                p.font.size = SLIDE_FONT_SIZE
        
        # Optionally auto-fit text into the textbox (if supported)
        try:
            text_frame.fit_text(max_size=SLIDE_FONT_SIZE.pt, min_size=12)
        except Exception as e:
            print("fit_text() not supported or failed:", e)

    prs.save(f"generated_ppt/{topic}_presentation.pptx")

def main():
    st.title("Text to PPT generation using Groq models")
    topic = st.text_input("Enter the topic you want to generate ppt on.")
    generate_button = st.button("Generate PPT")
    
    if generate_button and topic:
        st.info("Generating PPT presentation.... Please Wait..")
        slide_titles = generate_slide_titles(topic)
        filtered_slide_titles = [item for item in slide_titles if item.strip() != '']
        print("Slide titles:", filtered_slide_titles)
        
        slide_contents = [generate_slide_content(title, topic) for title in filtered_slide_titles]
        print("Slide contents:", slide_contents)
        
        create_presentation(topic, filtered_slide_titles, slide_contents)
        print("Presentation Generated Successfully!!")
        st.success("Presentation Generated Successfully!!")
        st.markdown(get_ppt_download_link(topic), unsafe_allow_html=True)


def get_ppt_download_link(topic):
    ppt_filename = f"generated_ppt/{topic}_presentation.pptx"

    with open(ppt_filename, "rb") as file:
        ppt_contents = file.read()

    b64_ppt = base64.b64encode(ppt_contents).decode()
    
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{ppt_filename}">Download the PowerPoint Presentation</a>'


if __name__ == "__main__":
    main()
